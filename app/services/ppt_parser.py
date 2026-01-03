from pptx import Presentation
from typing import List, Dict

def parse_ppt(file_path: str) -> List[Dict]:
    """
    Parses a PowerPoint file and extracts text from each slide.

    Args:
        file_path (str): The path to the PowerPoint file.
        
    Returns:
        List[Dict]: A list of dictionaries, each containing the slide number and its text content.]
    """

    presentation = Presentation(file_path)
    slides_content = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []

        # 1. Read slides shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text_parts.append(text)
        full_text = " ".join(slide_text_parts)

        slides_content.append({
            "slide_number": index,
            "text": full_text,
            "has_text": bool(full_text)
        })
    return slides_content

